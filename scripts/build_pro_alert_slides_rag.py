# Full RAG deck - saves to docs/Pro_Alert_Slides_RAG.pptx or %TEMP% if docs locked
# python scripts/build_pro_alert_slides_rag.py

from pathlib import Path
import os
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG = RGBColor(0x0A, 0x0F, 0x1A)
PRIMARY = RGBColor(0x3B, 0x82, 0xF6)
PRIMARY2 = RGBColor(0x60, 0xA5, 0xFA)
TEXT = RGBColor(0xE2, 0xE8, 0xF0)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
ROW_ALT = RGBColor(0x1E, 0x29, 0x3B)
GREEN_BG = RGBColor(0xBB, 0xF7, 0xD0)
AMBER_BG = RGBColor(0xFE, 0xF3, 0xC7)
RED_BG = RGBColor(0xFE, 0xE2, 0xE2)
HEADER_BG = RGBColor(0x1E, 0x40, 0xAF)
DARK = RGBColor(0x0F, 0x17, 0x2A)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
AMBER = RGBColor(0xD9, 0x77, 0x06)
RED = RGBColor(0xDC, 0x26, 0x26)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def rag_from_score(score):
    if score <= 3:
        return "LOW", GREEN_BG
    if score <= 6:
        return "MEDIUM", AMBER_BG
    return "CRITICAL", RED_BG


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def title_box(slide, text, top=Inches(0.28), sz=Pt(24)):
    b = slide.shapes.add_textbox(Inches(0.45), top, Inches(9.1), Inches(0.7))
    p = b.text_frame.paragraphs[0]
    p.text = text
    p.font.size = sz
    p.font.bold = True
    p.font.color.rgb = PRIMARY2


def cell_set(cell, txt, pt=Pt(7), bold=False, fg=TEXT, bg=None):
    cell.text = txt
    tf = cell.text_frame
    tf.word_wrap = True
    for m in (tf.margin_left, tf.margin_right, tf.margin_top, tf.margin_bottom):
        m = Pt(3)
    for p in tf.paragraphs:
        p.font.size = pt
        p.font.bold = bold
        p.font.color.rgb = fg
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
    blank = prs.slide_layouts[6]

    # Title
    s0 = prs.slides.add_slide(blank)
    set_bg(s0, BG)
    title_box(s0, "Pro-Alert", top=Inches(1.8), sz=Pt(34))
    t = s0.shapes.add_textbox(Inches(0.5), Inches(2.65), Inches(9), Inches(1.4)).text_frame.paragraphs[0]
    t.text = "IoT security platform\nRisk register 1–10 · RAG (red/amber/green)\nLow / Medium / Critical impact"
    t.font.size = Pt(17)
    t.font.color.rgb = TEXT

    # Methodology
    s1 = prs.slides.add_slide(blank)
    set_bg(s1, BG)
    title_box(s1, "Risk methodology & RAG scale")
    tb = s1.shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(9), Inches(5.2))
    tf = tb.text_frame
    tf.word_wrap = True
    blurb = [
        "SCORE 1–10: Combined likelihood × impact after controls. Higher = more attention.",
        "IMPACT: Low (limited) · Medium (service/user) · Critical (breach/legal/total loss).",
        "",
        "RAG:",
        "  GREEN  LOW      Score 1–3  — acceptable residual; monitor.",
        "  AMBER  MEDIUM   Score 4–6  — active controls; owner review.",
        "  RED    CRITICAL Score 7–10 — treat immediately; no go-live without mitigation.",
    ]
    for i, line in enumerate(blurb):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT
        if "GREEN" in line:
            p.font.color.rgb = GREEN
        if "AMBER" in line:
            p.font.color.rgb = AMBER
        if "RED" in line and "CRITICAL" in line:
            p.font.color.rgb = RED

    risks = [
        (1, "Auth", "JWT/secret compromise", "Med", "Critical", 9, "Default/leaked JWT; repo secrets → full takeover.", "Strong JWT_SECRET; no .env in git; HttpOnly; TTL; refresh rotation; pre-deploy gate.", "DevOps"),
        (2, "Auth", "Unverified access", "Med", "Critical", 8, "Session before verify → API/dashboard abuse.", "No JWT on signup; verify-then-login; deps reject unverified; revoke on verify.", "Dev"),
        (3, "Data", "MongoDB exposure", "Low", "Critical", 7, "No TLS/open cluster → data leak.", "TLS; auth; least privilege; backups encrypted.", "DevOps"),
        (4, "Comms", "SMTP abuse", "Med", "Medium", 5, "Forged/bulk mail.", "Rate limits; app passwords; server-side only.", "Dev"),
        (5, "Config", "Prod misconfig", "Med", "High", 7, "CORS *; HTTP; bad APP_BASE_URL.", "Allowlist CORS; HTTPS; HSTS; checklist.", "DevOps"),
        (6, "Supply", "Dependency CVEs", "Med", "Medium", 6, "RCE via lib/image.", "Pinned deps; audit; minimal base image.", "Dev"),
        (7, "Payments", "Webhook forgery", "Low", "High", 6, "Fake Stripe events.", "Webhook secret; signature verify.", "Dev"),
        (8, "Ops", "Backup/restore fail", "Med", "Critical", 8, "Data loss; no recovery.", "Continuous backup; restore drill; RTO/RPO.", "DevOps"),
        (9, "Avail", "DoS / limit bypass", "Med", "Medium", 5, "API exhaustion.", "Rate limits; scale-out; health separate.", "DevOps"),
        (10, "Human", "Insider/admin abuse", "Low", "Critical", 7, "Privileged misuse.", "RBAC; audit; MFA; access review.", "DevOps"),
    ]

    def table_slide(title, rows_slice):
        sl = prs.slides.add_slide(blank)
        set_bg(sl, BG)
        title_box(sl, title)
        n = len(rows_slice) + 1
        tbl = sl.shapes.add_table(n, 9, Inches(0.3), Inches(0.98), Inches(9.4), Inches(6.1)).table
        hdrs = ["#", "Cat", "Risk", "L", "I", "Score", "RAG", "Controls", "Owner"]
        for c, h in enumerate(hdrs):
            cell_set(tbl.cell(0, c), h, Pt(8), True, WHITE, HEADER_BG)
        for r, row in enumerate(rows_slice, 1):
            no, cat, title_r, L, impact, score, desc, ctrl, owner = row
            rag, rag_bg = rag_from_score(score)
            cell_set(tbl.cell(r, 0), str(no), bg=ROW_ALT if r % 2 else None)
            cell_set(tbl.cell(r, 1), cat, bg=ROW_ALT if r % 2 else None)
            cell_set(tbl.cell(r, 2), (title_r + " " + desc)[:95], bg=ROW_ALT if r % 2 else None)
            cell_set(tbl.cell(r, 3), L, bg=ROW_ALT if r % 2 else None)
            cell_set(tbl.cell(r, 4), impact, bg=ROW_ALT if r % 2 else None)
            cell_set(tbl.cell(r, 5), str(score), fg=DARK, bg=rag_bg)
            cell_set(tbl.cell(r, 6), rag, fg=DARK, bg=rag_bg)
            cell_set(tbl.cell(r, 7), ctrl[:100], bg=ROW_ALT if r % 2 else None)
            cell_set(tbl.cell(r, 8), owner, bg=ROW_ALT if r % 2 else None)

    table_slide("Risk register 1–5 (RAG colour-coded)", risks[:5])
    table_slide("Risk register 6–10 (RAG colour-coded)", risks[5:])

    # Performance slide
    sp = prs.slides.add_slide(blank)
    set_bg(sp, BG)
    title_box(sp, "Performance & security — in depth")
    bx = sp.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.8))
    tf = bx.text_frame
    tf.word_wrap = True
    blocks = [
        ("Confidentiality", ["JWT only after verify; refresh hashed; logout revokes all.", "Bcrypt; complexity; lockout; optional MFA.", "Secrets env-only; CORS allowlist."]),
        ("Integrity", ["Stripe webhooks signed; HTTPS/HSTS.", "Correct APP_BASE_URL on email links."]),
        ("Availability / performance", ["FastAPI + Motor async; indexes on hot paths.", "Stateless scale-out; WebSockets reduce polling.", "Rate limits on auth; retention avoids unbounded growth."]),
    ]
    first = True
    for h, bs in blocks:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = h
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = PRIMARY
        for b in bs:
            p = tf.add_paragraph()
            p.text = "  • " + b
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT

    # Deployment slide
    sd = prs.slides.add_slide(blank)
    set_bg(sd, BG)
    title_box(sd, "Deployment & maintenance — in depth")
    bx2 = sd.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.8))
    tf2 = bx2.text_frame
    tf2.word_wrap = True
    blocks2 = [
        ("Deployment", ["Docker Compose / Railway / Render; uvicorn 0.0.0.0 $PORT.", "security_gate before prod; TLS Mongo; no wildcard CORS."]),
        ("Config", [".env / platform vars single source; no secrets in layers.", "Rebuild on rotation."]),
        ("Maintenance", ["Backups + restore drill; RPO/RTO.", "Post-deploy smoke: auth, notifications, webhooks.", "Rollback via image tag; /api/health monitoring."]),
    ]
    first = True
    for h, bs in blocks2:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        p.text = h
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = PRIMARY
        for b in bs:
            p = tf2.add_paragraph()
            p.text = "  • " + b
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT

    docs = Path(__file__).resolve().parent.parent / "docs"
    out = docs / "Pro_Alert_Slides_RAG.pptx"
    try:
        prs.save(str(out))
    except PermissionError:
        out = Path(tempfile.gettempdir()) / "Pro_Alert_Slides_RAG.pptx"
        prs.save(str(out))
        print("docs locked — saved to", out)
    else:
        print("Saved:", out)


if __name__ == "__main__":
    main()
