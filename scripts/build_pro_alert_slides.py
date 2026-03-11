"""
Build Pro-Alert themed PowerPoint with risk register + two topic slides.
Run from project root:
  pip install python-pptx
  python scripts/build_pro_alert_slides.py
Output: docs/Pro_Alert_Slides.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Pro-Alert branding (from web/assets/app.css)
BG = RGBColor(0x0A, 0x0F, 0x1A)       # --bg deep navy
PRIMARY = RGBColor(0x3B, 0x82, 0xF6)  # --primary
PRIMARY2 = RGBColor(0x60, 0xA5, 0xFA)  # --primary-2
TEXT = RGBColor(0xE2, 0xE8, 0xF0)     # light slate
MUTED = RGBColor(0x94, 0xA3, 0xB8)


def set_slide_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text, top=Inches(0.35)):
    box = slide.shapes.add_textbox(Inches(0.5), top, Inches(9), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY2
    p.alignment = PP_ALIGN.LEFT


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(6.85), Inches(9), Inches(0.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = MUTED


def add_bullet_block(slide, left, top, width, height, title, bullets):
    if title:
        tb = slide.shapes.add_textbox(left, top, width, Inches(0.35))
        tfp = tb.text_frame.paragraphs[0]
        tfp.text = title
        tfp.font.size = Pt(16)
        tfp.font.bold = True
        tfp.font.color.rgb = PRIMARY
        top = top + Inches(0.4)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT
        p.space_before = Pt(4)
        p.level = 0


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # blank

    # --- Slide 0: Title (Pro-Alert branding) ---
    s0 = prs.slides.add_slide(blank)
    set_slide_background(s0, BG)
    add_title(s0, "Pro-Alert", top=Inches(2.0))
    sub = s0.shapes.add_textbox(Inches(0.5), Inches(2.95), Inches(9), Inches(1.2))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "IoT security monitoring platform\nRisk register · Performance & security · Deployment & maintenance"
    sp.font.size = Pt(20)
    sp.font.color.rgb = TEXT
    sp.alignment = PP_ALIGN.CENTER

    # --- Slide 1: Risk register table ---
    slide1 = prs.slides.add_slide(blank)
    set_slide_background(slide1, BG)
    add_title(slide1, "Final risk register – Pro-Alert")

    rows, cols = 9, 5
    left, top = Inches(0.45), Inches(1.25)
    width, height = Inches(9.1), Inches(5.1)
    table = slide1.shapes.add_table(rows, cols, left, top, width, height).table

    headers = ["ID", "Risk", "L", "I", "Mitigation"]
    data = [
        ["R1", "Weak/leaked JWT or secrets", "M", "C", "JWT_SECRET generated; no .env in git; security gate"],
        ["R2", "Access without verification", "M", "H", "Verify before login; no session on signup until verified"],
        ["R3", "MongoDB exposed / no TLS", "L–M", "C", "mongodb+srv + auth; backups"],
        ["R4", "SMTP / notification abuse", "M", "M", "Rate limits; server-side secrets only"],
        ["R5", "Bad prod config (CORS, HTTP)", "M", "H", "CORS_ORIGINS; FORCE_HTTPS; checklist"],
        ["R6", "Dependency vulnerabilities", "M", "M", "Pinned requirements; rebuild images"],
        ["R7", "Stripe webhook forgery", "L", "H", "Webhook secret; signature validation"],
        ["R8", "No backup / failed restore", "M", "H", "Atlas backups; restore drill"],
    ]

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = BG
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY

    for r, row in enumerate(data, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = TEXT
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)

    add_footer(
        slide1,
        "Residual risk accepted only where controls reduce impact · Critical items gated before go-live",
    )

    # --- Slide 2: Performance & security ---
    slide2 = prs.slides.add_slide(blank)
    set_slide_background(slide2, BG)
    add_title(slide2, "Artefact performance & security")
    add_bullet_block(
        slide2,
        Inches(0.5),
        Inches(1.2),
        Inches(4.4),
        Inches(2.6),
        "Security",
        [
            "JWT + refresh rotation; logout revokes refresh; optional MFA",
            "Zero-tolerance email verification before API/dashboard",
            "bcrypt; strong password policy; lockout after failed attempts",
            "Rate limits on signup, login, password reset",
            "HTTPS; strict CORS; pre-deploy security checklist",
        ],
    )
    add_bullet_block(
        slide2,
        Inches(5.1),
        Inches(1.2),
        Inches(4.4),
        Inches(2.6),
        "Performance",
        [
            "FastAPI + Motor (async MongoDB)",
            "Indexes on users.email, devices, alerts",
            "Stateless API – scale behind reverse proxy",
            "WebSockets for real-time dashboard",
        ],
    )
    takeaway = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.15), Inches(9), Inches(0.75)
    )
    takeaway.fill.solid()
    takeaway.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    takeaway.line.color.rgb = PRIMARY
    takeaway.line.width = Pt(2)
    tf = takeaway.text_frame
    tf.margin_left = tf.margin_right = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = "Takeaway: Layered security (verify → login → JWT); performance via async I/O and indexing."
    p.font.size = Pt(14)
    p.font.color.rgb = PRIMARY2
    p.alignment = PP_ALIGN.LEFT
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # --- Slide 3: Deployment & maintenance ---
    slide3 = prs.slides.add_slide(blank)
    set_slide_background(slide3, BG)
    add_title(slide3, "Artefact deployment & maintenance")
    add_bullet_block(
        slide3,
        Inches(0.5),
        Inches(1.2),
        Inches(4.4),
        Inches(2.8),
        "Deployment",
        [
            "Docker Compose: API + MongoDB 7; env_file for secrets",
            "Railway / Render: pip + uvicorn; /api/health",
            "Single .env – APP_BASE_URL, CORS, SMTP, Twilio, Stripe",
            "Pre-deploy: TLS DB, HTTPS, email links to live URL",
        ],
    )
    add_bullet_block(
        slide3,
        Inches(5.1),
        Inches(1.2),
        Inches(4.4),
        Inches(2.8),
        "Maintenance",
        [
            "Restart policies; log auth/email failures",
            "Backups; retention/cleanup jobs",
            "Post-deploy: test auth + notifications; rollback via image tag",
            "Monitor /api/health; product alerts for device offline",
        ],
    )
    takeaway3 = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.35), Inches(9), Inches(0.75)
    )
    takeaway3.fill.solid()
    takeaway3.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    takeaway3.line.color.rgb = PRIMARY
    takeaway3.line.width = Pt(2)
    tf3 = takeaway3.text_frame
    tf3.margin_left = tf3.margin_right = Inches(0.15)
    p3 = tf3.paragraphs[0]
    p3.text = (
        "Takeaway: Repeatable Docker + env deploy; checklist-driven maintenance with backups and health checks."
    )
    p3.font.size = Pt(13)
    p3.font.color.rgb = PRIMARY2
    tf3.vertical_anchor = MSO_ANCHOR.MIDDLE

    out = Path(__file__).resolve().parent.parent / "docs" / "Pro_Alert_Slides.pptx"
    prs.save(str(out))
    print(f"Saved: {out}")
    return out


if __name__ == "__main__":
    build_presentation()
